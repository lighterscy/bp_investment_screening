"""BP document parsing interfaces."""

from __future__ import annotations

from pathlib import Path
from typing import Protocol

from bp_investment_screening.schemas import DocumentPage, DocumentParseResult


class DocumentParser(Protocol):
    name: str

    def parse(self, path: str | Path) -> DocumentParseResult:
        """Parse a BP file into page-level text."""


class PlainTextParser:
    name = "plain_text"

    def parse(self, path: str | Path) -> DocumentParseResult:
        source_path = Path(path)
        text = source_path.read_text(encoding="utf-8")
        pages = [
            DocumentPage(
                page_number=1,
                text=text,
                source_path=str(source_path),
            )
        ]
        return DocumentParseResult(
            source_path=source_path,
            pages=pages,
            parser_name=self.name,
        )


class PdfParser:
    name = "pymupdf"

    def parse(self, path: str | Path) -> DocumentParseResult:
        try:
            import fitz
        except ImportError as exc:
            raise RuntimeError(
                "PDF parsing requires the optional dependency `pymupdf`. "
                "Install with `pip install '.[docs]'`."
            ) from exc

        source_path = Path(path)
        pages: list[DocumentPage] = []
        with fitz.open(source_path) as doc:
            for index, page in enumerate(doc, start=1):
                pages.append(
                    DocumentPage(
                        page_number=index,
                        text=page.get_text("text").strip(),
                        source_path=str(source_path),
                    )
                )
        return DocumentParseResult(
            source_path=source_path,
            pages=pages,
            parser_name=self.name,
        )


class PptxParser:
    name = "python_pptx"

    def parse(self, path: str | Path) -> DocumentParseResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError(
                "PPTX parsing requires the optional dependency `python-pptx`. "
                "Install with `pip install '.[docs]'`."
            ) from exc

        source_path = Path(path)
        presentation = Presentation(source_path)
        pages: list[DocumentPage] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
            pages.append(
                DocumentPage(
                    page_number=index,
                    text="\n".join(text for text in texts if text),
                    source_path=str(source_path),
                )
            )
        return DocumentParseResult(
            source_path=source_path,
            pages=pages,
            parser_name=self.name,
        )


class UnsupportedDocumentParser:
    name = "unsupported"

    def parse(self, path: str | Path) -> DocumentParseResult:
        source_path = Path(path)
        raise NotImplementedError(
            f"Document parsing for {source_path.suffix or '<no suffix>'} is not implemented yet."
        )


def choose_parser(path: str | Path) -> DocumentParser:
    suffix = Path(path).suffix.lower()
    if suffix in {".txt", ".md"}:
        return PlainTextParser()
    if suffix == ".pdf":
        return PdfParser()
    if suffix == ".pptx":
        return PptxParser()
    return UnsupportedDocumentParser()
